#!/usr/bin/env python3
import argparse
import json
from pathlib import Path

from pptx import Presentation


def read_notes(slide):
    try:
        if not slide.has_notes_slide:
            return ""
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            return (notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        return ""
    return ""


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    source = Path(args.input)
    output = Path(args.output)

    prs = Presentation(str(source))

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_blocks = []
        text_shapes = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = (shape.text or "").strip()
                if text:
                    normalized = text.replace("\n", " ").strip()
                    text_blocks.append(normalized)
                    text_shapes.append(
                        {
                            "text": normalized,
                            "left": int(getattr(shape, "left", 0) or 0),
                            "top": int(getattr(shape, "top", 0) or 0),
                            "width": int(getattr(shape, "width", 0) or 0),
                            "height": int(getattr(shape, "height", 0) or 0),
                        }
                    )

        notes_text = read_notes(slide)
        notes = [line.strip() for line in notes_text.splitlines() if line.strip()] if notes_text else []

        slides.append(
            {
                "page": idx,
                "textBlocks": text_blocks,
                "notes": notes,
                "textShapes": text_shapes,
            }
        )

    payload = {
        "slideCount": len(prs.slides),
        "slideWidth": int(getattr(prs, "slide_width", 0) or 0),
        "slideHeight": int(getattr(prs, "slide_height", 0) or 0),
        "slides": slides,
    }

    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


if __name__ == "__main__":
    main()
